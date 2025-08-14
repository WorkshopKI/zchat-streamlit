#!/usr/bin/env python3
"""
ChatBot v1.0 - Hauptanwendung
==========================

Eine funktionsreiche Chat-Anwendung mit KI-Unterstützung für Support-Ticket-Bearbeitung.

Hauptfunktionen:
- Projektbasierte Chat-Verwaltung mit Sessions
- Dokument-Upload und kontextbewusste KI-Antworten
- Mehrere KI-Provider (LM Studio, OpenRouter, Azure OpenAI)
- Persistent SQLite-Datenbankspeicherung
- Deutsche Benutzeroberfläche

Autor: ChatBot v1.0 Team
Version: 1.0
"""

import streamlit as st
import streamlit_antd_components as sac
import json
import uuid
from datetime import datetime
from typing import Dict, List, Optional
import base64
from pathlib import Path
from llm_integration import generate_context_aware_response, get_llm_provider
from config_manager import ConfigManager
from storage_service import get_storage_service, migrate_session_data

# Page config
st.set_page_config(
    page_title="Projekte",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Hide Streamlit's automatic navigation
st.markdown("""
<style>
/* Hide the navigation bar */
[data-testid="stSidebarNav"] {
    display: none !important;
}

/* Hide navigation for different Streamlit versions */
section[data-testid="stSidebarNav"] {
    display: none !important;
}

.css-1d391kg .css-1v0mbdj {
    display: none !important;
}
</style>
""", unsafe_allow_html=True)

# ============================================================================
# STYLING UND UI KONFIGURATION
# ============================================================================

def load_css(theme="light"):
    """
    Lädt vereinfachte CSS-Stile für Antd-Komponenten.
    
    Args:
        theme (str): Theme-Name ('light' oder 'dark')
        
    Note:
        Minimale CSS-Anpassungen, da Antd-Komponenten eigenes Styling mitbringen.
    """
    st.markdown("""
    <style>
    /* Minimal custom CSS for Antd integration */
    
    /* Wider sidebar default with user resizing preserved */
    .css-1d391kg {
        width: 330px;
        min-width: 280px;
        max-width: 500px;
    }
    
    .css-1cypcdb {
        width: 330px;
        min-width: 280px;
        max-width: 500px;
    }
    
    /* For newer Streamlit versions */
    section[data-testid="stSidebar"] {
        width: 330px;
        min-width: 280px;
        max-width: 500px;
    }
    
    /* Allow main content to adjust dynamically */
    .css-18e3th9 {
        margin-left: auto;
    }
    
    /* Override Streamlit's default button colors to neutral gray */
    .stButton > button {
        background-color: #f5f5f5 !important;
        color: #262730 !important;
        border: 1px solid #d9d9d9 !important;
        border-radius: 6px !important;
    }
    
    .stButton > button:hover {
        background-color: #e6f4ff !important;
        border-color: #91caff !important;
        color: #262730 !important;
    }
    
    .stButton > button:active, .stButton > button:focus {
        background-color: #bae0ff !important;
        border-color: #69b1ff !important;
        color: #262730 !important;
        box-shadow: none !important;
    }
    
    /* Token counter styling */
    .token-counter {
        background-color: #1890ff;
        color: white;
        padding: 4px 8px;
        border-radius: 4px;
        font-size: 0.8rem;
        margin-left: 8px;
        display: inline-block;
    }
    
    /* Chat message container improvements */
    .stChatMessage {
        padding: 1rem;
        margin: 0.5rem 0;
        border-radius: 8px;
        border: 1px solid #d9d9d9;
    }
    
    /* Ensure Antd components display properly */
    .ant-btn, .ant-tree, .ant-tabs, .ant-alert {
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif !important;
    }
    </style>
    """, unsafe_allow_html=True)

# ============================================================================
# ANWENDUNGSINITIALISIERUNG
# ============================================================================

def init_session_state():
    """
    Initialisiert den Streamlit Session State und die Datenbank.
    
    Führt folgende Aktionen aus:
    - Datenbank-Cache leeren für frische Instanz
    - Storage Service initialisieren
    - Konfiguration laden
    - Aktuelles Projekt aus Benutzereinstellungen laden
    - Einmalige Datenmigration durchführen
    """
    # Clear database cache to ensure fresh instance with new schema
    from database import clear_database_cache
    if 'cache_cleared' not in st.session_state:
        clear_database_cache()
        # Also clear storage service to get fresh database manager
        if 'storage_service' in st.session_state:
            del st.session_state.storage_service
        st.session_state.cache_cleared = True
    
    # Initialize storage service
    storage = get_storage_service()
    
    # Initialize config manager
    if 'config_manager' not in st.session_state:
        st.session_state.config_manager = ConfigManager()
    
    # Load settings from database or use defaults
    if 'settings' not in st.session_state:
        default_settings = st.session_state.config_manager.config
        loaded_settings = storage.load_settings(default_settings)
        st.session_state.settings = loaded_settings
        # Update config manager with loaded settings
        st.session_state.config_manager.config = loaded_settings
    
    # Initialize current project from user preferences
    if 'current_project' not in st.session_state:
        st.session_state.current_project = storage.get_user_preference('current_project')
    
    # Migration: Move existing session data to database (one-time operation)
    migrate_session_data()

# ============================================================================
# PROJEKT- UND CHAT-VERWALTUNG
# ============================================================================

def create_project(name: str, description: str = "") -> str:
    """
    Erstellt ein neues Projekt in der Datenbank.
    
    Args:
        name (str): Name des Projekts
        description (str): Optionale Beschreibung
        
    Returns:
        str: ID des erstellten Projekts
    """
    storage = get_storage_service()
    project_id = storage.create_project(name, description)
    # Save current project preference
    storage.save_user_preference('current_project', project_id)
    return project_id

def delete_project(project_id: str):
    """
    Löscht ein Projekt aus der Datenbank.
    
    Args:
        project_id (str): ID des zu löschenden Projekts
    """
    storage = get_storage_service()
    storage.delete_project(project_id)
    if st.session_state.current_project == project_id:
        st.session_state.current_project = None
        storage.save_user_preference('current_project', None)

def get_all_projects() -> List[Dict]:
    """
    Ruft alle Projekte aus der Datenbank ab.
    
    Returns:
        List[Dict]: Liste aller Projekte mit Metadaten
    """
    storage = get_storage_service()
    return storage.get_all_projects()

def get_chat_history(project_id: str, limit: int = None, session_id: str = None) -> List[Dict]:
    """
    Ruft den Chat-Verlauf für ein Projekt oder eine Session ab.
    
    Args:
        project_id (str): ID des Projekts
        limit (int, optional): Maximale Anzahl Nachrichten
        session_id (str, optional): Spezifische Session-ID
        
    Returns:
        List[Dict]: Liste der Chat-Nachrichten
    """
    storage = get_storage_service()
    return storage.get_chat_history(project_id, limit, session_id)

def add_message(project_id: str, role: str, content: str, avatar: Optional[str] = None, session_id: str = None):
    """
    Fügt eine Nachricht zum Chat-Verlauf hinzu.
    
    Args:
        project_id (str): ID des Projekts
        role (str): Rolle ('user' oder 'assistant')
        content (str): Nachrichteninhalt
        avatar (str, optional): Avatar-URL
        session_id (str, optional): Session-ID
    """
    storage = get_storage_service()
    metadata = {'avatar': avatar} if avatar else None
    storage.add_message(project_id, role, content, metadata, session_id)

# ============================================================================
# DOKUMENT-VERWALTUNG
# ============================================================================

def upload_document(project_id: str, filename: str, content: str, file_type: str = "text"):
    """
    Lädt ein Dokument zu einem Projekt hoch.
    
    Args:
        project_id (str): ID des Projekts
        filename (str): Name der Datei
        content (str): Dateiinhalt
        file_type (str): MIME-Type der Datei
        
    Returns:
        str: ID des hochgeladenen Dokuments
    """
    storage = get_storage_service()
    file_size = len(content.encode('utf-8'))
    doc_id = storage.add_document(project_id, filename, content, file_type, file_size)
    return doc_id

def get_documents(project_id: str) -> List[Dict]:
    """
    Ruft alle Dokumente eines Projekts ab.
    
    Args:
        project_id (str): ID des Projekts
        
    Returns:
        List[Dict]: Liste aller Projektdokumente
    """
    storage = get_storage_service()
    return storage.get_project_documents(project_id)

def get_document_content(project_id: str) -> Dict:
    """
    Ruft den Inhalt aller Dokumente eines Projekts ab.
    
    Args:
        project_id (str): ID des Projekts
        
    Returns:
        Dict: Dictionary mit Dokument-IDs als Keys und Inhalten als Values
    """
    storage = get_storage_service()
    return storage.get_document_content(project_id)

# ============================================================================
# HILFSFUNKTIONEN
# ============================================================================

def count_tokens(text: str) -> int:
    """
    Schätzt die Anzahl der Tokens in einem Text.
    
    Args:
        text (str): Zu zählender Text
        
    Returns:
        int: Geschätzte Token-Anzahl
        
    Note:
        Vereinfachte Schätzung: Wörter × 1.3
    """
    return len(text.split()) * 1.3  # Rough approximation

def extract_text_from_pdf(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Text aus PDF-Dateien mit OCR-Fallback.
    
    Args:
        file_content (bytes): PDF-Dateiinhalt
        filename (str): Name der Datei
        
    Returns:
        str: Extrahierter Text
    """
    import io
    import PyPDF2
    import fitz  # pymupdf
    import pytesseract
    from pdf2image import convert_from_bytes
    from PIL import Image
    
    text_content = ""
    
    try:
        # Method 1: Try PyPDF2 for text-based PDFs
        pdf_stream = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_stream)
        
        for page_num, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            if page_text.strip():
                text_content += f"\n## Seite {page_num + 1}\n\n{page_text}\n"
        
        # If we got substantial text, return it
        if len(text_content.strip()) > 100:
            return text_content.strip()
            
    except Exception as e:
        print(f"PyPDF2 extraction failed: {e}")
    
    try:
        # Method 2: Try PyMuPDF for better text extraction
        pdf_stream = io.BytesIO(file_content)
        pdf_document = fitz.open(stream=pdf_stream, filetype="pdf")
        
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            page_text = page.get_text()
            if page_text.strip():
                text_content += f"\n## Seite {page_num + 1}\n\n{page_text}\n"
        
        pdf_document.close()
        
        # If we got substantial text, return it
        if len(text_content.strip()) > 100:
            return text_content.strip()
            
    except Exception as e:
        print(f"PyMuPDF extraction failed: {e}")
    
    try:
        # Method 3: OCR fallback for image-based PDFs
        images = convert_from_bytes(file_content, dpi=200)
        
        for page_num, image in enumerate(images):
            # OCR with Tesseract (German + English)
            ocr_text = pytesseract.image_to_string(image, lang='deu+eng')
            if ocr_text.strip():
                text_content += f"\n## Seite {page_num + 1} (OCR)\n\n{ocr_text}\n"
        
        if text_content.strip():
            return text_content.strip()
            
    except Exception as e:
        print(f"OCR extraction failed: {e}")
    
    return f"**Fehler beim Extrahieren des PDF-Inhalts von {filename}**\n\nDie Datei konnte nicht verarbeitet werden. Möglicherweise ist sie passwortgeschützt oder beschädigt."

def extract_text_from_docx(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Text aus DOCX-Dateien.
    
    Args:
        file_content (bytes): DOCX-Dateiinhalt
        filename (str): Name der Datei
        
    Returns:
        str: Extrahierter und formatierter Text
    """
    import io
    from docx import Document
    
    try:
        doc_stream = io.BytesIO(file_content)
        doc = Document(doc_stream)
        
        content_parts = []
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                # Check if it's a heading based on style
                if paragraph.style.name.startswith('Heading'):
                    level = paragraph.style.name.replace('Heading ', '')
                    if level.isdigit():
                        content_parts.append(f"{'#' * int(level)} {text}")
                    else:
                        content_parts.append(f"## {text}")
                else:
                    content_parts.append(text)
        
        # Extract tables
        for table_num, table in enumerate(doc.tables):
            content_parts.append(f"\n### Tabelle {table_num + 1}\n")
            
            # Create markdown table
            for row_num, row in enumerate(table.rows):
                row_cells = [cell.text.strip() for cell in row.cells]
                if row_num == 0:
                    # Header row
                    content_parts.append("| " + " | ".join(row_cells) + " |")
                    content_parts.append("| " + " | ".join(["---"] * len(row_cells)) + " |")
                else:
                    content_parts.append("| " + " | ".join(row_cells) + " |")
        
        return "\n\n".join(content_parts) if content_parts else f"**Keine lesbaren Inhalte in {filename} gefunden.**"
        
    except Exception as e:
        return f"**Fehler beim Extrahieren des DOCX-Inhalts von {filename}**: {str(e)}"

def extract_text_from_xlsx(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Daten aus XLSX-Dateien und konvertiert zu Markdown.
    
    Args:
        file_content (bytes): XLSX-Dateiinhalt
        filename (str): Name der Datei
        
    Returns:
        str: Extrahierte Daten als Markdown
    """
    import io
    import pandas as pd
    
    try:
        excel_stream = io.BytesIO(file_content)
        
        # Read all sheets
        excel_file = pd.ExcelFile(excel_stream)
        content_parts = []
        
        for sheet_name in excel_file.sheet_names:
            try:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                if not df.empty:
                    content_parts.append(f"## Arbeitsblatt: {sheet_name}\n")
                    
                    # Convert to markdown table
                    markdown_table = df.to_markdown(index=False, tablefmt='pipe')
                    content_parts.append(markdown_table)
                    
                    # Add summary statistics for numeric columns
                    numeric_cols = df.select_dtypes(include=['number']).columns
                    if len(numeric_cols) > 0:
                        content_parts.append(f"\n### Zusammenfassung für {sheet_name}")
                        for col in numeric_cols:
                            stats = df[col].describe()
                            content_parts.append(f"**{col}**: Mittelwert={stats['mean']:.2f}, Min={stats['min']}, Max={stats['max']}")
                        
            except Exception as e:
                content_parts.append(f"**Fehler beim Lesen von Arbeitsblatt '{sheet_name}'**: {str(e)}")
        
        return "\n\n".join(content_parts) if content_parts else f"**Keine lesbaren Daten in {filename} gefunden.**"
        
    except Exception as e:
        return f"**Fehler beim Extrahieren der XLSX-Daten von {filename}**: {str(e)}"

def extract_text_from_pptx(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Text aus PPTX-Dateien.
    
    Args:
        file_content (bytes): PPTX-Dateiinhalt
        filename (str): Name der Datei
        
    Returns:
        str: Extrahierter und formatierter Text
    """
    import io
    from pptx import Presentation
    
    try:
        pptx_stream = io.BytesIO(file_content)
        presentation = Presentation(pptx_stream)
        
        content_parts = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            content_parts.append(f"## Folie {slide_num}")
            
            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                content_parts.append("\n".join(slide_text))
            else:
                content_parts.append("*[Keine Textinhalte gefunden]*")
            
            content_parts.append("---")  # Slide separator
        
        return "\n\n".join(content_parts) if content_parts else f"**Keine lesbaren Inhalte in {filename} gefunden.**"
        
    except Exception as e:
        return f"**Fehler beim Extrahieren der PPTX-Inhalte von {filename}**: {str(e)}"

def extract_text_from_image(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Text aus Bilddateien mit OCR.
    """
    import io
    import pytesseract
    from PIL import Image
    import cv2
    import numpy as np
    
    try:
        # Convert bytes to PIL Image
        image = Image.open(io.BytesIO(file_content))
        
        # Convert to RGB if necessary
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # OCR with German + English
        ocr_text = pytesseract.image_to_string(image, lang='deu+eng')
        
        if ocr_text.strip():
            return f"# OCR-Extraktion: {filename}\n\n{ocr_text.strip()}"
        else:
            return f"**Keine Textinhalte in {filename} erkannt.**\n\nDas Bild enthält möglicherweise keinen lesbaren Text oder die Qualität ist zu niedrig für OCR."
            
    except Exception as e:
        return f"**Fehler beim OCR-Prozess für {filename}**: {str(e)}"

def extract_text_from_html(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Text aus HTML-Dateien.
    """
    from bs4 import BeautifulSoup
    import chardet
    
    try:
        # Detect encoding
        detected = chardet.detect(file_content)
        encoding = detected['encoding'] or 'utf-8'
        
        html_content = file_content.decode(encoding)
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        content_parts = []
        
        # Extract title
        title = soup.find('title')
        if title:
            content_parts.append(f"# {title.get_text().strip()}")
        
        # Extract headings and content
        for tag in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'div', 'article', 'section']):
            text = tag.get_text().strip()
            if text:
                if tag.name.startswith('h'):
                    level = int(tag.name[1])
                    content_parts.append(f"{'#' * level} {text}")
                else:
                    content_parts.append(text)
        
        # Extract tables
        for table in soup.find_all('table'):
            rows = []
            for tr in table.find_all('tr'):
                cells = [td.get_text().strip() for td in tr.find_all(['td', 'th'])]
                if cells:
                    rows.append(cells)
            
            if rows:
                content_parts.append("\n### Tabelle\n")
                for i, row in enumerate(rows):
                    content_parts.append("| " + " | ".join(row) + " |")
                    if i == 0:  # Header separator
                        content_parts.append("| " + " | ".join(["---"] * len(row)) + " |")
        
        return "\n\n".join(content_parts) if content_parts else f"**Keine lesbaren Inhalte in {filename} gefunden.**"
        
    except Exception as e:
        return f"**Fehler beim Extrahieren der HTML-Inhalte von {filename}**: {str(e)}"

def extract_text_from_xml(file_content: bytes, filename: str) -> str:
    """
    Extrahiert strukturierte Daten aus XML-Dateien.
    """
    import xml.etree.ElementTree as ET
    import chardet
    
    try:
        # Detect encoding
        detected = chardet.detect(file_content)
        encoding = detected['encoding'] or 'utf-8'
        
        xml_content = file_content.decode(encoding)
        root = ET.fromstring(xml_content)
        
        content_parts = [f"# XML-Struktur: {filename}\n"]
        
        def process_element(element, level=0):
            indent = "  " * level
            parts = []
            
            # Element name and attributes
            if element.attrib:
                attrs = ", ".join([f"{k}={v}" for k, v in element.attrib.items()])
                parts.append(f"{indent}**{element.tag}** ({attrs})")
            else:
                parts.append(f"{indent}**{element.tag}**")
            
            # Element text
            if element.text and element.text.strip():
                parts.append(f"{indent}  {element.text.strip()}")
            
            # Process children
            for child in element:
                parts.extend(process_element(child, level + 1))
            
            return parts
        
        content_parts.extend(process_element(root))
        
        return "\n".join(content_parts)
        
    except Exception as e:
        return f"**Fehler beim Extrahieren der XML-Inhalte von {filename}**: {str(e)}"

def extract_text_from_rtf(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Text aus RTF-Dateien.
    """
    from striprtf.striprtf import rtf_to_text
    
    try:
        rtf_content = file_content.decode('utf-8')
        text_content = rtf_to_text(rtf_content)
        
        return f"# RTF-Dokument: {filename}\n\n{text_content}" if text_content.strip() else f"**Keine lesbaren Inhalte in {filename} gefunden.**"
        
    except Exception as e:
        return f"**Fehler beim Extrahieren der RTF-Inhalte von {filename}**: {str(e)}"

def extract_text_from_odt(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Text aus ODT-Dateien (OpenDocument Text).
    """
    import zipfile
    import io
    from xml.etree import ElementTree as ET
    
    try:
        # ODT is a ZIP file containing XML
        with zipfile.ZipFile(io.BytesIO(file_content)) as odt_zip:
            content_xml = odt_zip.read('content.xml')
            
        # Parse the content XML
        root = ET.fromstring(content_xml)
        
        # Define OpenDocument namespace
        namespaces = {
            'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
            'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0'
        }
        
        content_parts = []
        
        # Extract paragraphs and headings
        for para in root.findall('.//text:p', namespaces):
            text = ''.join(para.itertext()).strip()
            if text:
                content_parts.append(text)
        
        for heading in root.findall('.//text:h', namespaces):
            text = ''.join(heading.itertext()).strip()
            level = int(heading.get('{urn:oasis:names:tc:opendocument:xmlns:text:1.0}outline-level', '1'))
            if text:
                content_parts.append(f"{'#' * level} {text}")
        
        return "\n\n".join(content_parts) if content_parts else f"**Keine lesbaren Inhalte in {filename} gefunden.**"
        
    except Exception as e:
        return f"**Fehler beim Extrahieren der ODT-Inhalte von {filename}**: {str(e)}"

def extract_text_from_epub(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Text aus EPUB-Dateien.
    """
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    import io
    
    try:
        book = epub.read_epub(io.BytesIO(file_content))
        content_parts = []
        
        # Extract metadata
        title = book.get_metadata('DC', 'title')
        if title:
            content_parts.append(f"# {title[0][0]}")
        
        author = book.get_metadata('DC', 'creator')
        if author:
            content_parts.append(f"**Autor:** {author[0][0]}")
        
        content_parts.append("---\n")
        
        # Extract chapters
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                
                text = soup.get_text()
                if text.strip():
                    content_parts.append(f"## Kapitel: {item.get_name()}")
                    content_parts.append(text.strip())
        
        return "\n\n".join(content_parts) if content_parts else f"**Keine lesbaren Inhalte in {filename} gefunden.**"
        
    except Exception as e:
        return f"**Fehler beim Extrahieren der EPUB-Inhalte von {filename}**: {str(e)}"

def extract_text_from_email(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Text aus E-Mail-Dateien (EML, MSG).
    """
    import email
    from email.policy import default
    import chardet
    
    try:
        if filename.endswith('.msg'):
            # MSG format requires special handling
            import extract_msg
            msg = extract_msg.Message(io.BytesIO(file_content))
            
            content_parts = [
                f"# E-Mail: {filename}",
                f"**Von:** {msg.sender}",
                f"**An:** {msg.to}",
                f"**Betreff:** {msg.subject}",
                f"**Datum:** {msg.date}",
                "---",
                msg.body or "Kein Textinhalt verfügbar."
            ]
            
            # Extract attachments info
            if msg.attachments:
                content_parts.append("\n## Anhänge:")
                for attachment in msg.attachments:
                    content_parts.append(f"- {attachment.longFilename} ({attachment.size} Bytes)")
            
            return "\n\n".join(content_parts)
            
        else:
            # EML format
            detected = chardet.detect(file_content)
            encoding = detected['encoding'] or 'utf-8'
            
            email_content = file_content.decode(encoding)
            msg = email.message_from_string(email_content, policy=default)
            
            content_parts = [
                f"# E-Mail: {filename}",
                f"**Von:** {msg['From']}",
                f"**An:** {msg['To']}",
                f"**Betreff:** {msg['Subject']}",
                f"**Datum:** {msg['Date']}",
                "---"
            ]
            
            # Extract body
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        content_parts.append(part.get_content())
                        break
            else:
                content_parts.append(msg.get_content())
            
            return "\n\n".join(content_parts)
            
    except Exception as e:
        return f"**Fehler beim Extrahieren der E-Mail-Inhalte von {filename}**: {str(e)}"

def extract_text_from_archive(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Informationen aus Archiv-Dateien.
    """
    import zipfile
    import tarfile
    import py7zr
    import io
    
    try:
        content_parts = [f"# Archiv: {filename}", "## Enthaltene Dateien:"]
        
        if filename.endswith('.zip'):
            with zipfile.ZipFile(io.BytesIO(file_content)) as archive:
                for info in archive.filelist:
                    size = f"{info.file_size:,} Bytes" if info.file_size else "Ordner"
                    content_parts.append(f"- **{info.filename}** ({size})")
                    
        elif filename.endswith(('.tar', '.tar.gz', '.tgz')):
            with tarfile.open(fileobj=io.BytesIO(file_content)) as archive:
                for member in archive.getmembers():
                    size = f"{member.size:,} Bytes" if member.isfile() else "Ordner"
                    content_parts.append(f"- **{member.name}** ({size})")
                    
        elif filename.endswith('.7z'):
            with py7zr.SevenZipFile(io.BytesIO(file_content)) as archive:
                for info in archive.list():
                    size = f"{info.uncompressed:,} Bytes" if not info.is_directory else "Ordner"
                    content_parts.append(f"- **{info.filename}** ({size})")
        
        content_parts.append(f"\n**Gesamt:** {len(content_parts)-2} Einträge")
        
        return "\n".join(content_parts)
        
    except Exception as e:
        return f"**Fehler beim Extrahieren der Archiv-Inhalte von {filename}**: {str(e)}"

def extract_text_from_specialized(file_content: bytes, filename: str) -> str:
    """
    Extrahiert Text aus spezialisierten Formaten.
    """
    import json
    import yaml
    import configparser
    from icalendar import Calendar
    import vobject
    import chardet
    import io
    
    try:
        if filename.endswith('.json'):
            content = json.loads(file_content.decode('utf-8'))
            return f"# JSON-Datei: {filename}\n\n```json\n{json.dumps(content, indent=2, ensure_ascii=False)}\n```"
            
        elif filename.endswith(('.yaml', '.yml')):
            content = yaml.safe_load(file_content)
            return f"# YAML-Datei: {filename}\n\n```yaml\n{yaml.dump(content, default_flow_style=False, allow_unicode=True)}\n```"
            
        elif filename.endswith(('.ini', '.conf')):
            config = configparser.ConfigParser()
            config.read_string(file_content.decode('utf-8'))
            
            content_parts = [f"# Konfiguration: {filename}"]
            for section in config.sections():
                content_parts.append(f"\n## [{section}]")
                for key, value in config[section].items():
                    content_parts.append(f"**{key}:** {value}")
            
            return "\n".join(content_parts)
            
        elif filename.endswith('.ics'):
            cal = Calendar.from_ical(file_content)
            content_parts = [f"# Kalender: {filename}"]
            
            for event in cal.walk('vevent'):
                summary = event.get('summary', 'Kein Titel')
                start = event.get('dtstart')
                end = event.get('dtend')
                description = event.get('description', '')
                
                content_parts.append(f"\n## {summary}")
                if start:
                    content_parts.append(f"**Start:** {start.dt}")
                if end:
                    content_parts.append(f"**Ende:** {end.dt}")
                if description:
                    content_parts.append(f"**Beschreibung:** {description}")
            
            return "\n".join(content_parts)
            
        elif filename.endswith('.vcf'):
            vcard_text = file_content.decode('utf-8')
            content_parts = [f"# Kontakte: {filename}"]
            
            for vcard in vobject.readComponents(vcard_text):
                if hasattr(vcard, 'fn'):
                    content_parts.append(f"\n## {vcard.fn.value}")
                if hasattr(vcard, 'org'):
                    content_parts.append(f"**Organisation:** {vcard.org.value[0]}")
                if hasattr(vcard, 'email'):
                    content_parts.append(f"**E-Mail:** {vcard.email.value}")
                if hasattr(vcard, 'tel'):
                    content_parts.append(f"**Telefon:** {vcard.tel.value}")
            
            return "\n".join(content_parts)
            
        elif filename.endswith('.log'):
            detected = chardet.detect(file_content)
            encoding = detected['encoding'] or 'utf-8'
            log_content = file_content.decode(encoding)
            
            lines = log_content.split('\n')
            content_parts = [f"# Log-Datei: {filename}", f"**Zeilen:** {len(lines)}"]
            
            # Show first and last 20 lines
            if len(lines) > 40:
                content_parts.append("\n## Erste 20 Zeilen:")
                content_parts.append("```\n" + "\n".join(lines[:20]) + "\n```")
                content_parts.append("\n## Letzte 20 Zeilen:")
                content_parts.append("```\n" + "\n".join(lines[-20:]) + "\n```")
            else:
                content_parts.append("\n## Vollständiger Inhalt:")
                content_parts.append("```\n" + log_content + "\n```")
            
            return "\n".join(content_parts)
            
        elif filename.endswith('.tex'):
            from pylatexenc.latex2text import LatexNodes2Text
            
            latex_content = file_content.decode('utf-8')
            converter = LatexNodes2Text()
            text_content = converter.latex_to_text(latex_content)
            
            return f"# LaTeX-Dokument: {filename}\n\n{text_content}"
            
        elif filename.endswith('.ipynb'):
            notebook = json.loads(file_content.decode('utf-8'))
            content_parts = [f"# Jupyter Notebook: {filename}"]
            
            for i, cell in enumerate(notebook.get('cells', [])):
                cell_type = cell.get('cell_type', 'unknown')
                source = ''.join(cell.get('source', []))
                
                if cell_type == 'markdown':
                    content_parts.append(f"\n## Markdown-Zelle {i+1}")
                    content_parts.append(source)
                elif cell_type == 'code':
                    content_parts.append(f"\n## Code-Zelle {i+1}")
                    content_parts.append(f"```python\n{source}\n```")
                    
                    # Include outputs if available
                    outputs = cell.get('outputs', [])
                    if outputs:
                        content_parts.append("### Ausgabe:")
                        for output in outputs:
                            if 'text' in output:
                                content_parts.append("```\n" + ''.join(output['text']) + "\n```")
            
            return "\n".join(content_parts)
            
        elif filename.endswith(('.py', '.js', '.java', '.cpp', '.c', '.php', '.rb', '.go', '.rs')):
            detected = chardet.detect(file_content)
            encoding = detected['encoding'] or 'utf-8'
            code_content = file_content.decode(encoding)
            
            # Determine language for syntax highlighting
            lang_map = {
                '.py': 'python', '.js': 'javascript', '.java': 'java',
                '.cpp': 'cpp', '.c': 'c', '.php': 'php', '.rb': 'ruby',
                '.go': 'go', '.rs': 'rust'
            }
            
            ext = '.' + filename.split('.')[-1] if '.' in filename else ''
            lang = lang_map.get(ext, 'text')
            
            return f"# Code-Datei: {filename}\n\n```{lang}\n{code_content}\n```"
        
        else:
            return None  # Let main function handle it
            
    except Exception as e:
        return f"**Fehler beim Extrahieren der Inhalte von {filename}**: {str(e)}"

def convert_document_to_markdown(file_content: bytes, filename: str, file_type: str) -> str:
    """
    Konvertiert verschiedene Dokumenttypen zu Markdown.
    
    Args:
        file_content (bytes): Dateiinhalt als Bytes
        filename (str): Name der Datei
        file_type (str): MIME-Type der Datei
        
    Returns:
        str: Markdown-formatierter Inhalt
    """
    # Handle text files
    if filename.endswith('.md'):
        return file_content.decode('utf-8')
    elif filename.endswith('.txt') or file_type.startswith('text/'):
        content = file_content.decode('utf-8')
        return f"```text\n{content}\n```"
    
    # Handle CSV files specially
    elif filename.endswith('.csv'):
        import pandas as pd
        import io
        try:
            df = pd.read_csv(io.BytesIO(file_content))
            markdown_table = df.to_markdown(index=False, tablefmt='pipe')
            
            content_parts = [
                f"# CSV-Datei: {filename}",
                f"**Zeilen:** {len(df)}, **Spalten:** {len(df.columns)}",
                markdown_table
            ]
            
            # Add summary for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                content_parts.append("\n## Statistiken:")
                for col in numeric_cols:
                    stats = df[col].describe()
                    content_parts.append(f"**{col}**: Ø{stats['mean']:.2f}, Min={stats['min']}, Max={stats['max']}")
            
            return "\n\n".join(content_parts)
        except Exception as e:
            return f"**Fehler beim Verarbeiten der CSV-Datei {filename}**: {str(e)}"
    
    # Handle Office documents
    elif filename.endswith('.pdf') or file_type == 'application/pdf':
        return extract_text_from_pdf(file_content, filename)
    
    elif filename.endswith('.docx') or file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        return extract_text_from_docx(file_content, filename)
    
    elif filename.endswith('.xlsx') or file_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
        return extract_text_from_xlsx(file_content, filename)
    
    elif filename.endswith('.pptx') or file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        return extract_text_from_pptx(file_content, filename)
    
    # Handle images with OCR
    elif filename.endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif', '.webp')) or file_type.startswith('image/'):
        return extract_text_from_image(file_content, filename)
    
    # Handle web formats
    elif filename.endswith('.html') or filename.endswith('.htm') or file_type == 'text/html':
        return extract_text_from_html(file_content, filename)
    
    elif filename.endswith('.xml') or file_type == 'application/xml' or file_type == 'text/xml':
        return extract_text_from_xml(file_content, filename)
    
    # Handle rich text formats
    elif filename.endswith('.rtf'):
        return extract_text_from_rtf(file_content, filename)
    
    elif filename.endswith('.odt'):
        return extract_text_from_odt(file_content, filename)
    
    # Handle e-books
    elif filename.endswith('.epub'):
        return extract_text_from_epub(file_content, filename)
    
    # Handle email formats
    elif filename.endswith(('.eml', '.msg')):
        return extract_text_from_email(file_content, filename)
    
    # Handle archives
    elif filename.endswith(('.zip', '.rar', '.7z', '.tar', '.tar.gz', '.tgz')):
        return extract_text_from_archive(file_content, filename)
    
    # Handle specialized formats
    else:
        specialized_result = extract_text_from_specialized(file_content, filename)
        if specialized_result:
            return specialized_result
        
        # Final fallback for unknown types
        try:
            content = file_content.decode('utf-8')
            return f"```\n{content}\n```"
        except UnicodeDecodeError:
            return f"**📄 {filename}**\n\n> **Dateityp:** {file_type}\n> **Größe:** {len(file_content):,} Bytes\n> **Status:** Binärdatei - Inhalt kann nicht angezeigt werden\n\n*Diese Datei wurde hochgeladen, aber der Inhalt kann nicht als Text dargestellt werden.*"

def convert_to_markdown(content: str, filename: str, file_type: str) -> str:
    """
    Legacy function for backward compatibility with already processed documents.
    """
    if filename.endswith('.md'):
        return content
    elif filename.endswith('.txt') or 'text' in file_type:
        return f"```text\n{content}\n```"
    elif content.startswith("Binary file:"):
        return f"""**📄 {filename}**

> **Dateityp:** {file_type}
> **Status:** Veraltetes Format - neu hochladen für vollständige Konvertierung

**Aktuelle Vorschau:**
```
{content}
```

*Laden Sie diese Datei erneut hoch, um die vollständige Textextraktion zu erhalten.*"""
    else:
        return f"```\n{content}\n```"

def load_avatar(avatar_type: str) -> str:
    """
    Lädt Avatar-Bilder für die Chat-Anzeige.
    
    Args:
        avatar_type (str): Typ des Avatars ('assistant' oder 'user')
        
    Returns:
        str: Base64-kodiertes Bild oder None bei Fehler
    """
    try:
        if avatar_type == "assistant":
            with open("chat-logo-system.png", "rb") as f:
                return base64.b64encode(f.read()).decode()
        elif avatar_type == "user":
            with open("chat-logo-user.png", "rb") as f:
                return base64.b64encode(f.read()).decode()
    except FileNotFoundError:
        return None
    return None

# ============================================================================
# HAUPTANWENDUNG
# ============================================================================

def main():
    """
    Hauptfunktion der ChatBot-Anwendung.
    
    Verwaltet:
    - Seitenleiste mit Projektnavigation
    - Haupt-Chat-Interface mit Tabs
    - Dokument-Upload und -verwaltung
    - KI-Antwortgenerierung
    """
    init_session_state()
    load_css('light')  # Always use light theme
    
    # Sidebar for project management
    with st.sidebar:
        # Project selection with Antd Tree component
        
        # Get storage service for session operations
        storage = get_storage_service()
        
        # Migrate existing projects to sessions if needed
        if 'sessions_migrated' not in st.session_state:
            storage.migrate_existing_projects_to_sessions()
            st.session_state.sessions_migrated = True
        
        # Project list with Antd Tree view in scrollable container
        projects = get_all_projects()
        if projects:
            # Sort projects alphabetically by name
            projects_sorted = sorted(projects, key=lambda proj: proj['name'].lower())
            
            # Initialize current session if not set
            if 'current_session' not in st.session_state:
                st.session_state.current_session = None
            
            # Build tree structure for Antd Tree component using TreeItem objects
            tree_items = []
            current_project_id = st.session_state.get('current_project')
            current_session_id = st.session_state.get('current_session')
            
            for project in projects_sorted:
                project_id = project['id']
                project_name = project['name']
                message_count = project['message_count']
                
                # Check if this project is currently selected
                is_selected_project = project_id == current_project_id
                
                # Get sessions for this project
                sessions = storage.get_project_sessions(project_id)
                session_count = len(sessions)
                
                # Build session children for tree using simplified TreeItem objects
                session_children = []
                for session in sessions:
                    is_selected_session = (session['id'] == current_session_id and is_selected_project)
                    
                    # Simple text-based indicators - arrow for active sessions
                    session_indicator = " ▶" if is_selected_session else ""
                    session_text = f"💬 {session['name']} ({session['message_count']}){session_indicator}"
                    
                    session_children.append(
                        sac.TreeItem(session_text)  # Minimal TreeItem structure
                    )
                
                # Simple text-based indicators for projects
                project_emoji = '📂' if is_selected_project else '📁'
                project_indicator = " ✓" if is_selected_project else ""
                project_text = f"{project_emoji} {project_name} ({session_count} Sessions){project_indicator}"
                
                # Add project to tree using simplified TreeItem objects
                tree_items.append(
                    sac.TreeItem(
                        project_text,
                        children=session_children
                    )
                )
            
            # Create scrollable container for tree
            with st.container():
                st.markdown('<div style="max-height: 600px; overflow-y: auto; border: 1px solid #e0e0e0; border-radius: 4px; padding: 8px;">', unsafe_allow_html=True)
                
                # Render Antd Tree component with fresh key to clear any cached parameters
                tree_key = f"project_tree_clean_v3_{len(projects_sorted)}"
                selected_items = sac.tree(
                    items=tree_items,
                    label='Projekt wählen:',
                    format_func='title',
                    open_all=False,
                    open_index=[0],  # Expand first project
                    return_index=True,
                    size='sm',
                    color='gray',
                    key=tree_key
                )
                
                st.markdown('</div>', unsafe_allow_html=True)
            
            # Handle tree selection with change detection to prevent infinite reruns
            if selected_items is not None:
                try:
                    # Store the last processed selection to detect changes
                    last_selection_key = 'last_tree_selection'
                    current_selection_value = str(selected_items)
                    
                    # Only process if selection actually changed
                    if st.session_state.get(last_selection_key) != current_selection_value:
                        st.session_state[last_selection_key] = current_selection_value
                        
                        # Handle both integer and list returns from tree component
                        if isinstance(selected_items, (list, tuple)):
                            if len(selected_items) > 0:
                                selected_index = selected_items[0]  # Get first selected index
                            else:
                                selected_index = None
                        else:
                            selected_index = selected_items  # Direct integer index
                        
                        # Parse the selection based on tree structure only if we have a valid index
                        if selected_index is not None and isinstance(selected_index, int):
                            # Build flat mapping of indices to project/session IDs
                            flat_mapping = []
                            for i, project in enumerate(projects_sorted):
                                flat_mapping.append(('project', project['id']))
                                sessions = storage.get_project_sessions(project['id'])
                                for session in sessions:
                                    flat_mapping.append(('session', session['id'], project['id']))
                            
                            if 0 <= selected_index < len(flat_mapping):
                                selection = flat_mapping[selected_index]
                                selection_changed = False
                                
                                if selection[0] == 'project':
                                    project_id = selection[1]
                                    if st.session_state.current_project != project_id:
                                        st.session_state.current_project = project_id
                                        storage.save_user_preference('current_project', project_id)
                                        # Auto-select first session
                                        sessions = storage.get_project_sessions(project_id)
                                        if sessions:
                                            st.session_state.current_session = sessions[0]['id']
                                        selection_changed = True
                                        
                                elif selection[0] == 'session':
                                    session_id = selection[1]
                                    project_id = selection[2]
                                    if (st.session_state.current_session != session_id or 
                                        st.session_state.current_project != project_id):
                                        st.session_state.current_project = project_id
                                        st.session_state.current_session = session_id
                                        storage.save_user_preference('current_project', project_id)
                                        selection_changed = True
                                
                                # Only rerun if selection actually changed
                                if selection_changed:
                                    st.rerun()
                
                except Exception as e:
                    st.error(f"Error handling tree selection: {str(e)}")
                    # Reset to prevent further errors
                    st.session_state.current_project = None
                    st.session_state.current_session = None
            
        else:
            sac.alert(
                label="Noch keine Projekte vorhanden",
                description="Gehen Sie zur Einstellungsseite, um Ihr erstes Projekt zu erstellen!",
                banner=False,
                closable=False,
                key='no_projects_alert'
            )
        
        # New Chat Creation Section - Only show when a project is selected
        if (projects and 
            st.session_state.get('current_project') and 
            st.session_state.current_project is not None):
            with st.expander("➕ Neuer Chat", expanded=False):
                with st.container():
                    # Input field with helpful placeholder
                    new_chat_name = st.text_input(
                        "Chat Name",
                        placeholder="z.B. Support Ticket #123, Bug Report, Kundenanfrage...",
                        key="new_chat_input",
                        label_visibility="collapsed",
                        help="Geben Sie einen beschreibenden Namen für den neuen Chat ein"
                    )
                    
                    # Action buttons with live preview
                    col1, col2 = st.columns([1, 2])
                    
                    with col1:
                        create_btn = sac.buttons([
                            sac.ButtonsItem(
                                label='Erstellen',
                                icon='plus-circle',
                                color='#52c41a'
                            )
                        ], 
                        variant='filled',
                        size='sm',
                        key='create_chat_btn',
                        return_index=True
                        )
                    
                    with col2:
                        # Live preview of chat name
                        if new_chat_name.strip():
                            preview_text = new_chat_name[:25] + ('...' if len(new_chat_name) > 25 else '')
                            st.markdown(f"**Vorschau:** {preview_text}")
                    
                    # Handle chat creation with Antd alerts
                    if create_btn is not None:
                        if new_chat_name.strip():
                            try:
                                # Create new chat session
                                session_id = storage.create_chat_session(
                                    st.session_state.current_project, 
                                    new_chat_name.strip()
                                )
                                if session_id:
                                    # Switch to new session
                                    st.session_state.current_session = session_id
                                    sac.alert(
                                        label=f"✅ Chat '{new_chat_name}' erfolgreich erstellt!",
                                        banner=False,
                                        closable=True,
                                        key='chat_created_alert'
                                    )
                                    # Clear input field
                                    st.session_state.new_chat_input = ""
                                    st.rerun()
                                else:
                                    sac.alert(
                                        label="❌ Fehler beim Erstellen des Chats",
                                        banner=False,
                                        closable=True,
                                        key='chat_error_alert'
                                    )
                            except Exception as e:
                                sac.alert(
                                    label=f"❌ Fehler: {str(e)}",
                                    banner=False,
                                    closable=True,
                                    key='chat_exception_alert'
                                )
                        else:
                            sac.alert(
                                label="⚠️ Bitte einen Chat-Namen eingeben",
                                banner=False,
                                closable=True,
                                key='chat_name_required_alert'
                            )
        
        sac.divider(key='nav_divider1')
        
        # Navigation to Settings page - using standard button for reliability
        # Reset return flag if it exists
        if 'just_returned_from_settings' not in st.session_state:
            st.session_state.just_returned_from_settings = False
            
        # Only show navigation if we haven't just returned
        if not st.session_state.just_returned_from_settings:
            if st.button("⚙️ Zu den Einstellungen", key='settings_nav_reliable', use_container_width=True):
                st.switch_page("pages/Settings.py")
        else:
            # Show disabled button when just returned, then reset flag
            st.button("⚙️ Zu den Einstellungen", key='settings_nav_disabled', use_container_width=True, disabled=True)
            st.session_state.just_returned_from_settings = False
        
        sac.divider(key='nav_divider2')
        
        # LLM Selection (moved below create project)
        with st.expander("⚙️ LLM Auswahl"):
            config_manager = st.session_state.config_manager
            
            # Provider selection
            enabled_providers = config_manager.get_enabled_providers()
            provider_names = list(enabled_providers.keys())
            
            if not provider_names:
                st.error("Keine LLM-Provider sind aktiviert!")
                st.info("Gehen Sie zur Einstellungsseite, um Provider zu konfigurieren")
            else:
                current_provider = st.session_state.settings.get('model_provider', provider_names[0])
                if current_provider not in provider_names:
                    current_provider = provider_names[0]
                
                new_provider = st.selectbox(
                    "Aktiver Provider", 
                    provider_names,
                    index=provider_names.index(current_provider),
                    format_func=lambda x: enabled_providers[x].get('name', x.title())
                )
                
                if new_provider != current_provider:
                    st.session_state.settings['model_provider'] = new_provider
                    config_manager.set('llm_providers.default_provider', new_provider)
                    # Save to database
                    storage = get_storage_service()
                    storage.save_settings(st.session_state.settings)
        
    
    # Main content area
    if not st.session_state.current_project:
        st.title("Willkommen bei ChatBot v1.0")
        st.info("Bitte erstellen oder wählen Sie ein Projekt aus der Seitenleiste, um zu beginnen.")
        return
    
    # Get current project info from database
    storage = get_storage_service()
    current_project = storage.get_project(st.session_state.current_project)
    
    # Ensure we have a current session for this project
    if not st.session_state.current_session:
        sessions = storage.get_project_sessions(st.session_state.current_project)
        if sessions:
            st.session_state.current_session = sessions[0]['id']
        else:
            # Create default session
            st.session_state.current_session = storage.create_chat_session(
                st.session_state.current_project, "Haupt-Chat"
            )
    
    # Project header with current session info
    col1, col2, col3 = st.columns([3, 1, 1])
    with col1:
        # Show project name and current session
        session_info = ""
        if st.session_state.current_session:
            sessions = storage.get_project_sessions(st.session_state.current_project)
            current_session = next((s for s in sessions if s['id'] == st.session_state.current_session), None)
            if current_session:
                session_info = f" - {current_session['name']}"
        
        st.title(f"{current_project['name']}{session_info}")
    with col2:
        # Token counter for current session
        session_messages = get_chat_history(st.session_state.current_project, session_id=st.session_state.current_session)
        total_tokens = sum(count_tokens(msg['content']) for msg in session_messages)
        st.markdown(f"""
        <div class="token-counter">
            Token: {int(total_tokens/1000)}k von 62k
        </div>
        """, unsafe_allow_html=True)
    with col3:
        # Compact export dropdown
        export_option = st.selectbox(
            "Export",
            ["", "📥 Session exportieren"],
            key="export_dropdown",
            label_visibility="collapsed",
            placeholder="Export..."
        )
        
        # Handle export selection
        if export_option == "📥 Session exportieren":
            # Export current session
            session_messages = get_chat_history(st.session_state.current_project, session_id=st.session_state.current_session)
            current_session = next((s for s in storage.get_project_sessions(st.session_state.current_project) 
                                  if s['id'] == st.session_state.current_session), None)
            export_data = {
                'project': current_project,
                'session': current_session,
                'chat_history': session_messages,
                'exported_at': datetime.now().isoformat()
            }
            session_name = current_session['name'] if current_session else 'session'
            st.download_button(
                label="📥 Session-Export herunterladen",
                data=json.dumps(export_data, indent=2),
                file_name=f"session_export_{session_name}.json",
                mime="application/json",
                use_container_width=True
            )
            # Reset dropdown after selection
            st.session_state.export_dropdown = ""
    
    # Antd Tabs for Chat and Documents
    doc_count = len(get_documents(st.session_state.current_project))
    active_tab = sac.tabs([
        sac.TabsItem(label='Chat', icon='message'),
        sac.TabsItem(label=f'Dokumente ({doc_count})', icon='file-text'),
    ], align='left', return_index=True, key='main_tabs')
    
    if active_tab == 0:
        # Chat interface
        chat_container = st.container()
        
        # Display chat history for current session
        with chat_container:
            session_messages = get_chat_history(st.session_state.current_project, session_id=st.session_state.current_session)
            
            # Welcome message if no chat history in this session
            if not session_messages:
                with st.chat_message("assistant", avatar="🤖"):
                    st.markdown("Willkommen! Ich bin Dein Assistent und unterstütze Dich bei der Bearbeitung.")
                    st.markdown("**Support-Workflow:**")
                    
                    # Show support workflow with Antd Steps
                    sac.steps(
                        items=[
                            sac.StepsItem(title='First Step', subtitle='Erste Analyse'),
                            sac.StepsItem(title='Second Step', subtitle='Details'),
                            sac.StepsItem(title='Third Step', subtitle='Spezialist'),
                        ],
                        index=0,
                        format_func='title',
                        # size='lg',
                        color='red',
                        placement='vertical',
                        return_index=True,
                        key='support_workflow_steps'
                    )
                    
                    st.markdown("Mit welcher Phase möchtest Du beginnen?")
            
            # Display existing messages from current session
            for message in session_messages:
                avatar = "🤖" if message['role'] == 'assistant' else "👤"
                with st.chat_message(message['role'], avatar=avatar):
                    st.markdown(message['content'])
        
        # Chat input
        if prompt := st.chat_input("Nachricht eingeben..."):
            # Add user message to current session
            add_message(st.session_state.current_project, "user", prompt, session_id=st.session_state.current_session)
            
            # Display user message
            with st.chat_message("user", avatar="👤"):
                st.markdown(prompt)
            
            # Generate assistant response
            with st.chat_message("assistant", avatar="🤖"):
                # Prepare messages for LLM (get session-specific history)
                chat_messages = []
                session_history = get_chat_history(st.session_state.current_project, session_id=st.session_state.current_session)
                for msg in session_history:
                    chat_messages.append({
                        "role": msg["role"],
                        "content": msg["content"]
                    })
                
                # Add current user message
                chat_messages.append({"role": "user", "content": prompt})
                
                # Get documents for context (project-wide)
                documents = get_document_content(st.session_state.current_project)
                
                # Generate streaming response
                response_placeholder = st.empty()
                full_response = ""
                
                try:
                    provider_name = st.session_state.settings.get('model_provider', 'lm_studio')
                    for chunk in generate_context_aware_response(
                        chat_messages, 
                        documents, 
                        provider_name=provider_name,
                        config_manager=st.session_state.config_manager,
                        stream=True
                    ):
                        full_response += chunk
                        response_placeholder.markdown(full_response + "▌")
                    
                    # Final response without cursor
                    response_placeholder.markdown(full_response)
                    
                    # Add assistant message to current session
                    add_message(st.session_state.current_project, "assistant", full_response, session_id=st.session_state.current_session)
                    
                except Exception as e:
                    error_message = f"Entschuldigung, ich habe einen Fehler festgestellt: {str(e)}"
                    response_placeholder.markdown(error_message)
                    add_message(st.session_state.current_project, "assistant", error_message, session_id=st.session_state.current_session)
    
    elif active_tab == 1:
        # Documents tab
        st.subheader("Dokumentenverwaltung")
        
        col1, col2 = st.columns([3, 1])
        with col1:
            # Add custom CSS to increase file uploader height
            st.markdown("""
            <style>
            .stFileUploader > div > div > div > div {
                min-height: 200px !important;
                padding: 20px !important;
            }
            .stFileUploader > div > div > div > div > div {
                min-height: 160px !important;
                display: flex !important;
                align-items: center !important;
                justify-content: center !important;
            }
            </style>
            """, unsafe_allow_html=True)
            
            # Get max file size from config for dynamic help text
            config_manager = st.session_state.config_manager
            max_file_size_mb = config_manager.get("file_upload.max_file_size_mb", 50)
            
            uploaded_files = st.file_uploader(
                "Dateien hier ablegen oder zum Auswählen klicken",
                accept_multiple_files=True,
                type=['pdf', 'xlsx', 'docx', 'pptx', 'txt', 'md', 'csv', 'png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif', 'webp', 'html', 'htm', 'xml', 'rtf', 'odt', 'ods', 'odp', 'epub', 'eml', 'msg', 'zip', 'rar', '7z', 'tar', 'json', 'yaml', 'yml', 'ini', 'conf', 'ics', 'vcf', 'log', 'tex', 'ipynb', 'py', 'js', 'java', 'cpp', 'c', 'php', 'rb', 'go', 'rs'],
                help=f"Unterstützte Formate: PDF, Office-Dokumente, Bilder (mit OCR), Web-Formate, E-Books, E-Mails, Archive, Code-Dateien und viele mehr • Maximale Dateigröße: {max_file_size_mb}MB pro Datei • Mehrere Dateien gleichzeitig möglich"
            )
        with col2:
            # Add vertical spacing to align with file uploader label
            st.markdown("<div style='margin-top: 53px;'></div>", unsafe_allow_html=True)
            # st.markdown("**Hinweis:**", unsafe_allow_html=True)
            st.info("<- Dokumente hier hochladen. Viele Formate werden unterstützt.")
        
        # Process uploaded files with duplicate prevention and progress bar
        if uploaded_files:
            # Initialize upload tracking
            if 'processed_files' not in st.session_state:
                st.session_state.processed_files = set()
            
            # Get current documents to check for duplicates
            storage = get_storage_service()
            existing_docs = storage.get_project_documents(st.session_state.current_project)
            existing_filenames = {doc['filename'] for doc in existing_docs}
            
            # Filter out files that need processing
            files_to_process = []
            max_file_size_bytes = max_file_size_mb * 1024 * 1024
            
            for uploaded_file in uploaded_files:
                if uploaded_file is not None:
                    file_id = f"{uploaded_file.name}_{uploaded_file.size}_{st.session_state.current_project}"
                    
                    # Skip if already processed in this session
                    if file_id in st.session_state.processed_files:
                        continue
                    
                    # Check file size limit
                    if uploaded_file.size > max_file_size_bytes:
                        file_size_mb = uploaded_file.size / (1024 * 1024)
                        st.error(f"❌ Datei '{uploaded_file.name}' ist zu groß ({file_size_mb:.1f}MB). Maximale Größe: {max_file_size_mb}MB")
                        continue
                    
                    # Check if file already exists in database
                    if uploaded_file.name in existing_filenames:
                        st.warning(f"Datei '{uploaded_file.name}' existiert bereits und wurde übersprungen.")
                        continue
                    
                    files_to_process.append(uploaded_file)
            
            # Process files with progress bar
            if files_to_process:
                st.info(f"Verarbeite {len(files_to_process)} Datei(en)...")
                
                # Create progress bar
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                total_files = len(files_to_process)
                
                for i, uploaded_file in enumerate(files_to_process):
                    try:
                        # Update progress - Reading phase
                        progress = (i * 2) / (total_files * 2)  # Each file has 2 phases
                        progress_bar.progress(progress)
                        status_text.text(f"📖 Lese Datei: {uploaded_file.name}")
                        
                        # Create unique file identifier
                        file_id = f"{uploaded_file.name}_{uploaded_file.size}_{st.session_state.current_project}"
                        
                        # Read file content as bytes
                        file_bytes = uploaded_file.read()
                        
                        # Update progress - Converting/Uploading phase
                        progress = (i * 2 + 1) / (total_files * 2)
                        progress_bar.progress(progress)
                        status_text.text(f"💾 Konvertiere zu Markdown: {uploaded_file.name}")
                        
                        # Convert to markdown using comprehensive extraction
                        file_type = uploaded_file.type or "text/plain"
                        markdown_content = convert_document_to_markdown(file_bytes, uploaded_file.name, file_type)
                        
                        # Upload document with converted content
                        doc_id = upload_document(
                            st.session_state.current_project,
                            uploaded_file.name,
                            markdown_content,
                            file_type
                        )
                        
                        # Mark as processed
                        st.session_state.processed_files.add(file_id)
                        
                    except Exception as e:
                        st.error(f"Fehler beim Hochladen von '{uploaded_file.name}': {str(e)}")
                        continue
                
                # Complete progress
                progress_bar.progress(1.0)
                status_text.text(f"✅ Fertig! {total_files} Datei(en) erfolgreich hochgeladen.")
                
                # Show success summary
                st.success(f"Alle {total_files} Dateien wurden erfolgreich verarbeitet:")
                for uploaded_file in files_to_process:
                    st.write(f"   • {uploaded_file.name}")
                
                # Clear progress indicators after a moment
                import time
                time.sleep(1)
                progress_bar.empty()
                status_text.empty()
        
        # Display documents with pagination
        documents = get_documents(st.session_state.current_project)
        if documents:
            st.subheader("Hochgeladene Dokumente")
            
            # Add refresh button to clear any caching issues
            col_docs, col_refresh = st.columns([4, 1])
            with col_refresh:
                if st.button("🔄 Aktualisieren", key="refresh_docs", help="Dokumentliste neu laden"):
                    # Force refresh by clearing any cached data
                    if hasattr(st.session_state, 'storage_service'):
                        del st.session_state.storage_service
                    # Clear processed files tracking
                    if 'processed_files' in st.session_state:
                        del st.session_state.processed_files
                    st.rerun()
            
            with col_docs:
                st.write(f"**{len(documents)} Dokument(e) gefunden:**")
            
            # Pagination logic - show at least 5 documents
            docs_per_page = 5
            total_docs = len(documents)
            
            if total_docs <= docs_per_page:
                # Show all documents if 5 or fewer
                docs_to_show = documents
                show_pagination = False
            else:
                # Initialize pagination
                if 'doc_page' not in st.session_state:
                    st.session_state.doc_page = 0
                
                show_pagination = True
                start_idx = st.session_state.doc_page * docs_per_page
                end_idx = min(start_idx + docs_per_page, total_docs)
                docs_to_show = documents[start_idx:end_idx]
                
                # Pagination controls
                col_prev, col_info, col_next = st.columns([1, 2, 1])
                with col_prev:
                    if st.button("← Zurück", disabled=st.session_state.doc_page == 0):
                        st.session_state.doc_page -= 1
                        st.rerun()
                with col_info:
                    current_page = st.session_state.doc_page + 1
                    total_pages = (total_docs - 1) // docs_per_page + 1
                    st.write(f"Seite {current_page} von {total_pages}")
                with col_next:
                    if st.button("Weiter →", disabled=end_idx >= total_docs):
                        st.session_state.doc_page += 1
                        st.rerun()
            
            # Display documents in compact format
            for doc in docs_to_show:
                metadata = doc.get('metadata', {})
                char_count = metadata.get('char_count', len(doc.get('content', '')))
                
                # Compact document row
                col_info, col_del = st.columns([6, 1])
                
                with col_del:
                    if st.button("🗑️", key=f"del_{doc['id']}", help=f"Dokument '{doc['filename']}' löschen"):
                        storage = get_storage_service()
                        try:
                            storage.delete_document(doc['id'])
                            st.success(f"Dokument '{doc['filename']}' wurde gelöscht.")
                            # Clear storage service cache to ensure fresh data
                            if hasattr(st.session_state, 'storage_service'):
                                del st.session_state.storage_service
                            st.rerun()
                        except Exception as e:
                            st.error(f"Fehler beim Löschen: {str(e)}")
                
                with col_info:
                    # Compact expandable section
                    with st.expander(f"📄 {doc['filename']} • {char_count:,} Zeichen • {doc['upload_date'][:10]}"):
                        # Compact metadata in one line
                        st.caption(f"📋 {doc['file_type']} • 📏 {doc['file_size']:,} Bytes • 📅 {doc['upload_date'][:19]}")
                        
                        # Convert and display content as markdown
                        markdown_content = convert_to_markdown(
                            doc['content'], 
                            doc['filename'], 
                            doc['file_type']
                        )
                        st.markdown(markdown_content)
            
            # Bottom pagination (if needed)
            if show_pagination:
                st.write("---")
                col_prev2, col_info2, col_next2 = st.columns([1, 2, 1])
                with col_prev2:
                    if st.button("← Zurück", key="prev_bottom", disabled=st.session_state.doc_page == 0):
                        st.session_state.doc_page -= 1
                        st.rerun()
                with col_info2:
                    st.write(f"Seite {current_page} von {total_pages}")
                with col_next2:
                    if st.button("Weiter →", key="next_bottom", disabled=end_idx >= total_docs):
                        st.session_state.doc_page += 1
                        st.rerun()
        else:
            st.info("Noch keine Dokumente hochgeladen.")

if __name__ == "__main__":
    main()